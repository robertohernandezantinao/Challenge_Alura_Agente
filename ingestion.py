# ingestion.py
"""
Módulo de ingestión de documentos multi-formato para el
Agente de Conocimiento Corporativo.

Formatos soportados: PDF, DOCX, XLSX/XLS, CSV, PPTX, Markdown, JSON, HTML y TXT.

Flujo:
    1. load_documents(carpeta)   -> lee todos los archivos soportados y los
                                     convierte en objetos Document de LangChain.
    2. build_vectorstore(docs)   -> los divide en fragmentos (chunks), genera
                                     embeddings con Gemini y crea un índice FAISS.
    3. load_vectorstore()        -> recarga el índice ya creado para hacer
                                     búsquedas semánticas (usado por agent.py).
"""
import os
import json
import time
import pandas as pd
from bs4 import BeautifulSoup
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_community.vectorstores import FAISS

SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".xlsx", ".xls", ".csv",
    ".pptx", ".md", ".markdown", ".json", ".html", ".htm", ".txt",
}


# --------------------------------------------------------------------------
# Loaders individuales por tipo de archivo
# --------------------------------------------------------------------------
def _load_pdf(path):
    from pypdf import PdfReader
    reader = PdfReader(path)
    docs = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            docs.append(Document(page_content=text,
                                  metadata={"source": os.path.basename(path), "page": i + 1}))
    return docs


def _load_docx(path):
    import docx
    document = docx.Document(path)
    text = "\n".join(p.text for p in document.paragraphs if p.text.strip())
    return [Document(page_content=text, metadata={"source": os.path.basename(path)})]


def _load_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    docs = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        content = "\n".join(t for t in texts if t.strip())
        if content.strip():
            docs.append(Document(page_content=content,
                                  metadata={"source": os.path.basename(path), "slide": i + 1}))
    return docs


def _load_excel(path):
    xls = pd.ExcelFile(path)
    docs = []
    for sheet_name in xls.sheet_names:
        df = xls.parse(sheet_name)
        content = df.to_markdown(index=False)
        docs.append(Document(page_content=content,
                              metadata={"source": os.path.basename(path), "sheet": sheet_name}))
    return docs


def _load_csv(path):
    df = pd.read_csv(path)
    content = df.to_markdown(index=False)
    return [Document(page_content=content, metadata={"source": os.path.basename(path)})]


def _load_json(path):
    with open(path, "r", encoding="utf-8") as f:
        data = json.load(f)
    content = json.dumps(data, ensure_ascii=False, indent=2)
    return [Document(page_content=content, metadata={"source": os.path.basename(path)})]


def _load_html(path):
    with open(path, "r", encoding="utf-8") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    text = soup.get_text(separator="\n")
    text = "\n".join(line.strip() for line in text.splitlines() if line.strip())
    return [Document(page_content=text, metadata={"source": os.path.basename(path)})]


def _load_text(path):
    with open(path, "r", encoding="utf-8") as f:
        text = f.read()
    return [Document(page_content=text, metadata={"source": os.path.basename(path)})]


LOADERS = {
    ".pdf": _load_pdf,
    ".docx": _load_docx,
    ".pptx": _load_pptx,
    ".xlsx": _load_excel,
    ".xls": _load_excel,
    ".csv": _load_csv,
    ".json": _load_json,
    ".html": _load_html,
    ".htm": _load_html,
    ".md": _load_text,
    ".markdown": _load_text,
    ".txt": _load_text,
}


def load_documents(folder_path: str):
    """Recorre una carpeta (recursivamente, incluyendo subcarpetas por categoría, p. ej.
    docs/company/, docs/product/, docs/ai/, etc.) y carga todos los documentos soportados.

    A cada fragmento se le asigna:
      - metadata["source"] -> ruta relativa a folder_path (ej. "product/Roadmap del Producto.md"),
        para no perder el contexto de categoría cuando hay nombres de archivo repetidos.
      - metadata["area"]   -> la primera carpeta bajo folder_path (ej. "product", "ai", "hr"),
        útil para filtrar o etiquetar respuestas por dominio en el futuro.
    """
    all_docs = []
    if not os.path.isdir(folder_path):
        print(f"⚠️  La carpeta '{folder_path}' no existe, se omite.")
        return all_docs

    for root, _, files in os.walk(folder_path):
        for filename in files:
            ext = os.path.splitext(filename)[1].lower()
            if ext not in SUPPORTED_EXTENSIONS:
                continue
            full_path = os.path.join(root, filename)
            loader_fn = LOADERS.get(ext)
            try:
                docs = loader_fn(full_path)
                rel_path = os.path.relpath(full_path, folder_path).replace(os.sep, "/")
                area = rel_path.split("/")[0] if "/" in rel_path else "general"
                for d in docs:
                    d.metadata["source"] = rel_path
                    d.metadata["area"] = area
                all_docs.extend(docs)
                print(f"✅ Cargado: {rel_path} ({len(docs)} fragmento(s))")
            except Exception as e:
                print(f"⚠️  No se pudo cargar {filename}: {e}")
    return all_docs


def load_all_documents(folders=("docs", "sample_docs")):
    """Combina varias carpetas raíz en un solo corpus (ej. 'docs/' con la documentación
    curada de la empresa + 'sample_docs/' con lo que se suba manualmente desde la interfaz).
    Carpetas que no existan se ignoran sin error, así se puede escalar de 40 a 97 documentos
    simplemente agregando archivos dentro de las mismas subcarpetas de 'docs/'."""
    all_docs = []
    for folder in folders:
        all_docs.extend(load_documents(folder))
    return all_docs


def build_vectorstore(docs, persist_path="vectorstore_index", chunk_size=1000, chunk_overlap=150,
                       batch_size=80, delay_seconds=61):
    """Divide los documentos en chunks, genera embeddings EN LOTES (con pausas entre cada
    lote) y crea/guarda el índice FAISS.

    El tier gratuito de la API de Gemini limita las solicitudes de embeddings a ~100 por
    minuto. Con corpus grandes (decenas de documentos), enviar todos los chunks de una sola
    vez dispara un error 429 RESOURCE_EXHAUSTED. Por eso aquí se procesan de a `batch_size`
    fragmentos (por defecto 80, con margen respecto al límite de 100) y se espera
    `delay_seconds` (por defecto 61, un poco más de un minuto) entre lotes para que la cuota
    se reinicie antes de continuar.
    """
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    chunks = splitter.split_documents(docs)
    total = len(chunks)

    embeddings = GoogleGenerativeAIEmbeddings(model="models/gemini-embedding-001")

    vectorstore = None
    for i in range(0, total, batch_size):
        batch = chunks[i:i + batch_size]
        lote_num = i // batch_size + 1
        total_lotes = (total + batch_size - 1) // batch_size
        print(f"🧠 Generando embeddings — lote {lote_num}/{total_lotes} "
              f"(fragmentos {i + 1}–{min(i + batch_size, total)} de {total})")

        if vectorstore is None:
            vectorstore = FAISS.from_documents(batch, embeddings)
        else:
            vectorstore.add_documents(batch)

        # Solo esperar si todavía quedan lotes por procesar.
        if i + batch_size < total:
            print(f"⏳ Esperando {delay_seconds}s para no exceder la cuota gratuita de la API...")
            time.sleep(delay_seconds)

    vectorstore.save_local(persist_path)
    print(f"📦 Índice vectorial guardado en '{persist_path}' ({total} fragmentos)")
    return vectorstore


def load_vectorstore(persist_path="vectorstore_index"):
    """Carga un índice FAISS ya creado previamente."""
    embeddings = GoogleGenerativeAIEmbeddings(model="models/gemini-embedding-001")
    return FAISS.load_local(persist_path, embeddings, allow_dangerous_deserialization=True)


if __name__ == "__main__":
    # Ejecutar `python ingestion.py` procesa 'docs/' (documentación curada de la empresa)
    # y 'sample_docs/' (archivos de prueba/sueltos) y crea el índice combinado.
    documentos = load_all_documents(["docs", "sample_docs"])
    if documentos:
        build_vectorstore(documentos)
    else:
        print("No se encontraron documentos para procesar en 'sample_docs/'.")
